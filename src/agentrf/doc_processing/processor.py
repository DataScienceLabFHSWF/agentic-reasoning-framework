# processes documents into markdown text, with metadata and provenance info

from __future__ import annotations
import csv
import json
import logging
import re
from pathlib import Path
from typing import Callable, Dict, List, Optional

from tqdm import tqdm

from .models import ProcessedDocument

logger = logging.getLogger(__name__)

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    import docx
except Exception:
    docx = None

try:
    from markdownify import markdownify as md
except Exception:
    md = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


class DocProcessor:
    """
    Normalize supported document formats into markdown text.

    Main use cases:
    1. Process a single document
    2. Process an entire knowledge base folder

    Best practice:
    - return both text and processed_path
    - save processed markdown to disk for downstream reproducibility
    """

    def __init__(self, output_dir: str | Path | None = None) -> None:
        self.output_dir = Path(output_dir) if output_dir else None

        self._dispatch: Dict[str, Callable[[Path], Optional[str]]] = {
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".doc": self._process_docx,  # best-effort only, old .doc may fail
            ".html": self._process_html,
            ".htm": self._process_html,
            ".md": self._process_md,
            ".markdown": self._process_md,
            ".txt": self._process_text,
            ".csv": self._process_csv,
            ".json": self._process_json,
            ".pptx": self._process_pptx,
            ".ppt": self._process_pptx,  # best-effort only
        }

    def process_document(
        self,
        file_path: str | Path,
        output_dir: str | Path | None = None,
        save_output: bool = True,
    ) -> Optional[ProcessedDocument]:
        file_path = Path(file_path)

        if not file_path.is_file():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = file_path.suffix.lower()
        handler = self._dispatch.get(ext)
        if handler is None:
            logger.warning("Unsupported file type: %s", file_path)
            return None

        text = handler(file_path)
        if not text:
            logger.warning("No text extracted from file: %s", file_path)
            return None

        final_output_dir = Path(output_dir) if output_dir else self.output_dir
        processed_path: Optional[str] = None

        if save_output and final_output_dir is not None:
            final_output_dir.mkdir(parents=True, exist_ok=True)
            dest = self._build_output_path(file_path=file_path, output_dir=final_output_dir)
            dest.write_text(text, encoding="utf-8")
            processed_path = str(dest.resolve())

        return ProcessedDocument(
            text=text,
            source_path=str(file_path.resolve()),
            file_type=ext,
            processed_path=processed_path,
            metadata={
                "source_name": file_path.name,
                "source_stem": file_path.stem,
            },
        )

    def process_knowledge_base(
        self,
        input_dir: str | Path,
        output_dir: str | Path | None = None,
        save_output: bool = True,
    ) -> List[ProcessedDocument]:
        input_dir = Path(input_dir)

        if not input_dir.exists():
            raise FileNotFoundError(f"Input directory not found: {input_dir}")

        if not input_dir.is_dir():
            raise NotADirectoryError(f"Expected a directory, got: {input_dir}")

        all_files = [f for f in input_dir.rglob("*") if f.is_file() and f.suffix.lower() in self._dispatch]

        processed_docs: List[ProcessedDocument] = []

        for file_path in tqdm(all_files, desc="Processing knowledge base"):
            try:
                doc = self.process_document(
                    file_path=file_path,
                    output_dir=output_dir,
                    save_output=save_output,
                )
                if doc is not None:
                    processed_docs.append(doc)
            except Exception as exc:
                logger.exception("Failed to process %s: %s", file_path, exc)

        return processed_docs

    def _build_output_path(self, file_path: Path, output_dir: Path) -> Path:
        """
        Save as file.md instead of file.pdf.md.
        Avoid collisions by appending _1, _2, ...
        """
        dest = output_dir / f"{file_path.stem}.md"
        counter = 1

        while dest.exists():
            dest = output_dir / f"{file_path.stem}_{counter}.md"
            counter += 1

        return dest

    def _process_pdf(self, file_path: Path) -> Optional[str]:
        if fitz is None:
            logger.warning("PyMuPDF not installed. Cannot process PDF: %s", file_path)
            return None

        try:
            doc = fitz.open(str(file_path))
        except Exception:
            logger.exception("Could not open PDF: %s", file_path)
            return None

        text_parts: List[str] = []
        try:
            for page in doc:
                try:
                    page_text = page.get_text("text")
                    if page_text:
                        text_parts.append(page_text)
                except Exception:
                    logger.warning("Failed to read a page in PDF: %s", file_path)
                    continue
        finally:
            doc.close()

        return "\n\n".join(text_parts).strip()

    def _process_docx(self, file_path: Path) -> Optional[str]:
        if docx is None:
            logger.warning("python-docx not installed. Cannot process DOCX: %s", file_path)
            return None

        try:
            document = docx.Document(str(file_path))
        except Exception:
            logger.exception("Could not open DOCX/DOC: %s", file_path)
            return None

        lines: List[str] = []

        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name.lower() if para.style and para.style.name else ""

            if "heading 1" in style_name:
                lines.append(f"# {text}")
            elif "heading 2" in style_name:
                lines.append(f"## {text}")
            elif "heading 3" in style_name:
                lines.append(f"### {text}")
            else:
                lines.append(text)

        return "\n\n".join(lines).strip()

    def _process_html(self, file_path: Path) -> Optional[str]:
        try:
            raw = file_path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            logger.exception("Could not read HTML file: %s", file_path)
            return None

        if md is not None:
            try:
                return md(raw).strip()
            except Exception:
                logger.warning("markdownify failed for HTML: %s", file_path)

        return re.sub(r"<[^>]+>", "", raw).strip()

    def _process_md(self, file_path: Path) -> Optional[str]:
        try:
            return file_path.read_text(encoding="utf-8", errors="ignore").strip()
        except Exception:
            logger.exception("Could not read Markdown file: %s", file_path)
            return None

    def _process_text(self, file_path: Path) -> Optional[str]:
        try:
            return file_path.read_text(encoding="utf-8", errors="ignore").strip()
        except Exception:
            logger.exception("Could not read text file: %s", file_path)
            return None

    def _process_csv(self, file_path: Path) -> Optional[str]:
        try:
            with open(file_path, "r", encoding="utf-8", newline="") as f:
                reader = csv.reader(f)
                rows = list(reader)
        except Exception:
            logger.exception("Could not read CSV file: %s", file_path)
            return None

        if not rows:
            return ""

        header = rows[0]
        body = rows[1:]

        table_lines: List[str] = []
        table_lines.append("|" + "|".join(self._escape_md_cell(h) for h in header) + "|")
        table_lines.append("|" + "|".join(["---"] * len(header)) + "|")

        for row in body:
            padded = row + [""] * (len(header) - len(row))
            padded = padded[: len(header)]
            table_lines.append("|" + "|".join(self._escape_md_cell(cell) for cell in padded) + "|")

        return "\n".join(table_lines).strip()

    def _process_json(self, file_path: Path) -> Optional[str]:
        try:
            data = json.loads(file_path.read_text(encoding="utf-8"))
        except Exception:
            logger.exception("Could not read JSON file: %s", file_path)
            return None

        pretty = json.dumps(data, indent=2, ensure_ascii=False)
        return f"```json\n{pretty}\n```"

    def _process_pptx(self, file_path: Path) -> Optional[str]:
        if Presentation is None:
            logger.warning("python-pptx not installed. Cannot process PPTX: %s", file_path)
            return None

        try:
            pres = Presentation(str(file_path))
        except Exception:
            logger.exception("Could not open PPT/PPTX: %s", file_path)
            return None

        lines: List[str] = []

        for i, slide in enumerate(pres.slides, start=1):
            slide_lines: List[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(getattr(shape, "text", "")).strip()
                    if text:
                        slide_lines.append(text)

            if slide_lines:
                lines.append(f"## Slide {i}")
                lines.extend(slide_lines)
                lines.append("")

        return "\n".join(lines).strip()

    @staticmethod
    def _escape_md_cell(value: str) -> str:
        return (value or " ").replace("|", "\\|").replace("\n", " ").strip() or " "
