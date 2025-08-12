"""
doc_processor.py
----------------

This module defines a ``DocProcessor`` class that can ingest files from a
directory, normalise their contents into a common textual representation and
persist that representation to disk.  The processed text can then be
consumed by retrieval‑augmented generation (RAG) pipelines built on
frameworks such as LlamaIndex, LangChain or LangGraph.  By abstracting
away the specifics of each file type, this processor makes it easy to
ingest heterogeneous data sets (PDFs, Word documents, HTML pages, CSVs,
presentations, etc.) into a single corpus suitable for embedding or
keyword‑based retrieval.

The processor performs the following high‑level tasks:

* Walk a directory tree and discover files based on their extension.
* For each supported type, use an appropriate parser to extract the
  meaningful text.  For example, PyMuPDF is used for PDFs, python‑docx
  for Word files, markdownify for HTML, the csv module for CSVs and
  python‑pptx for PowerPoint files.
* Optionally add lightweight structure to the output (e.g. converting
  headings in a docx file into Markdown ``#`` prefixes) to aid downstream
  chunking.  If a converter cannot be imported (for example because
  optional dependencies are not installed), the processor falls back to a
  simple plain‑text extraction.
* Write the extracted text to a file under the ``processed_files``
  directory, preserving the original base name but changing the suffix to
  ``.md``.  Plain text files are still written with the ``.md`` suffix so
  that Markdown‑aware node parsers will treat them uniformly.
* Return a list of ``ProcessedDocument`` instances, each of which holds
  the extracted text and basic metadata about the source.

The design tries to strike a balance between preserving document
structure (where possible) and keeping the code simple and free from
heavy dependencies.  You can extend the processor by adding more
``_process_XXX`` methods and registering them in the ``_dispatch``
dictionary.

Usage example:

    >>> from doc_processor import DocProcessor
    >>> processor = DocProcessor()
    >>> docs = processor.process_directory("/path/to/raw", "./processed_files")
    >>> # docs is a list of ProcessedDocument instances
    >>> first_doc = docs[0]
    >>> print(first_doc.text[:500])

Note: To integrate the output with LlamaIndex or LangChain, you can
convert each ``ProcessedDocument`` into the corresponding framework’s
Document type.  For LlamaIndex, you would call ``Document(text=doc.text,
metadata=doc.metadata)``.  For LangChain, use
``langchain.schema.Document(page_content=doc.text, metadata=doc.metadata)``.
"""

from __future__ import annotations

import csv
import json
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Dict, List, Optional
from tqdm import tqdm

# Optional imports.  We attempt to import these libraries but do not
# require them for the processor to operate.  If any import fails we
# silently fall back to simple text extraction.
try:
    import fitz  # type: ignore[import]
except Exception:  # pragma: no cover
    fitz = None  # type: ignore

try:
    import docx  # type: ignore[import]
except Exception:  # pragma: no cover
    docx = None  # type: ignore

try:
    from markdownify import markdownify as md  # type: ignore[import]
except Exception:  # pragma: no cover
    md = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore[import]
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore


@dataclass
class ProcessedDocument:
    """A simple container for processed text and metadata.

    Attributes
    ----------
    text:
        The extracted textual content of the document.
    metadata:
        A dictionary containing at least a ``source`` key that records
        the original file path.  You can add arbitrary additional
        metadata in subclasses or downstream code.
    """

    text: str
    metadata: Dict[str, str] = field(default_factory=dict)


class DocProcessor:
    """Process a collection of files into a unified textual format.

    The processor walks through a given input directory, delegates file
    parsing based on the file extension and writes the processed output
    into an output directory.  Each processed file is written with a
    ``.md`` extension irrespective of the original extension.  This
    uniform naming scheme makes it easy to feed the documents into
    Markdown‑aware node parsers (e.g. ``MarkdownNodeParser`` in
    LlamaIndex) or treat them as plain text.
    """

    def __init__(self) -> None:
        # Map file suffixes to handler methods.  You can extend this
        # dictionary to support additional formats.  Each handler takes
        # ``file_path`` and returns a tuple of (text, output_suffix).
        self._dispatch: Dict[str, Callable[[Path], Optional[str]]] = {
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".doc": self._process_docx,
            ".html": self._process_html,
            ".htm": self._process_html,
            ".md": self._process_md,
            ".markdown": self._process_md,
            ".txt": self._process_text,
            ".csv": self._process_csv,
            ".json": self._process_json,
            ".pptx": self._process_pptx,
            ".ppt": self._process_pptx,
        }

    # Public API
    from tqdm import tqdm

    def process_directory(self, input_dir: str, output_dir: str) -> List[ProcessedDocument]:
        in_path = Path(input_dir)
        out_path = Path(output_dir)

        if not out_path.exists():
            out_path.mkdir(parents=True, exist_ok=True)
            print(f"📁 Created output directory: {out_path}")
        else:
            print(f"📁 Using existing output directory: {out_path}")
    
        processed_docs: List[ProcessedDocument] = []

        all_files = list(in_path.rglob("*"))
        supported_files = [f for f in all_files if f.is_file() and f.suffix.lower() in self._dispatch]

        print(f"📄 Found {len(supported_files)} supported files in: {input_dir}")
        current_dir = None

        for file_path in tqdm(supported_files, desc="📦 Processing documents"):
            if current_dir != file_path.parent:
                current_dir = file_path.parent
                print(f"\n📂 Processing folder: {current_dir}")

            ext = file_path.suffix.lower()
            handler = self._dispatch.get(ext)
            if handler is None:
                print(f"[⏭️ Skipped] Unsupported file: {file_path}")
                continue

            try:
                text = handler(file_path)
                if not text:
                    print(f"[⚠️ Skipped] Empty or unprocessable file: {file_path}")
                    continue
            except Exception as exc:
                print(f"[❌ ERROR] Failed to process {file_path}: {exc}")
                continue

            output_name = f"{file_path.stem}.md"
            dest = out_path / output_name
            try:
                with open(dest, "w", encoding="utf-8") as f:
                    f.write(text)
                print(f"[✅ Saved] → {dest.name}")
            except Exception as exc:
                print(f"[❌ ERROR] Failed to write {dest}: {exc}")
                continue

            processed_docs.append(
                ProcessedDocument(
                    text=text,
                    metadata={"source": str(file_path.resolve())},
                )
            )

        print(f"\n🎉 Done! Successfully processed {len(processed_docs)} files.")
        return processed_docs

    # Handler implementations

    def _process_pdf(self, file_path: Path) -> Optional[str]:
        """Extract text from a PDF using PyMuPDF.

        If PyMuPDF is not available, return ``None``.
        """
        if fitz is None:
            print(f"[Warning] fitz (PyMuPDF) is not installed; skipping PDF {file_path}")
            return None
        text_parts: List[str] = []
        try:
            doc = fitz.open(str(file_path))  # type: ignore[arg-type]
        except Exception as exc:
            print(f"[Warning] Failed to open PDF {file_path}: {exc}")
            return None
        for page in doc:
            try:
                # Extract plain text.  The 'text' option preserves basic layout
                # ordering without including coordinates.
                page_text = page.get_text("text")
                text_parts.append(page_text)
            except Exception as exc:
                print(f"[Warning] Failed to extract page {page.number} from {file_path}: {exc}")
                continue
        doc.close()
        return "\n\n".join(text_parts).strip()

    def _process_docx(self, file_path: Path) -> Optional[str]:
        """Extract text from a Word document.

        Attempts to preserve headings by converting them into Markdown
        headings.  If ``python-docx`` is not available, falls back to
        returning raw text by reading the file as binary and decoding it
        using utf-8 (which is almost certainly wrong for binary formats, so
        failure should be expected).
        """
        if docx is None:
            print(f"[Warning] python-docx is not installed; attempting naive read of {file_path}")
            try:
                return file_path.read_text(encoding="utf-8")
            except Exception:
                return None
        try:
            doc = docx.Document(str(file_path))  # type: ignore[arg-type]
        except Exception as exc:
            print(f"[Warning] Failed to open DOCX {file_path}: {exc}")
            return None
        lines: List[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = getattr(para.style, "name", "")
            # Convert Word heading styles into Markdown headings.  The
            # built‑in styles usually follow the pattern "Heading 1",
            # "Heading 2", etc.  If the style name does not start with
            # "Heading", we leave the text unchanged.
            if style_name.startswith("Heading"):
                try:
                    level_str = style_name.split(" ")[1]
                    level = int(level_str)
                except (IndexError, ValueError):
                    level = 0
                prefix = "#" * max(1, level)
                lines.append(f"{prefix} {text}")
            else:
                lines.append(text)
        return "\n".join(lines).strip()

    def _process_html(self, file_path: Path) -> Optional[str]:
        """Convert HTML into Markdown if markdownify is available.

        Otherwise falls back to stripping tags via a naive replacement.
        """
        try:
            raw = file_path.read_text(encoding="utf-8")
        except Exception as exc:
            print(f"[Warning] Failed to read HTML {file_path}: {exc}")
            return None
        if md is not None:
            try:
                return md(raw, heading_style="ATX")
            except Exception as exc:
                print(f"[Warning] markdownify failed on {file_path}: {exc}")
                # Fall back to naive text extraction
        # Simple fall back: strip out tags and collapse whitespace
        import re  # local import to avoid polluting module namespace
        text = re.sub(r"<[^>]+>", "", raw)
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    def _process_md(self, file_path: Path) -> Optional[str]:
        """Read Markdown/MD files as text."""
        try:
            return file_path.read_text(encoding="utf-8").strip()
        except Exception as exc:
            print(f"[Warning] Failed to read Markdown {file_path}: {exc}")
            return None

    def _process_text(self, file_path: Path) -> Optional[str]:
        """Read plain text files."""
        try:
            return file_path.read_text(encoding="utf-8").strip()
        except Exception as exc:
            print(f"[Warning] Failed to read text file {file_path}: {exc}")
            return None

    def _process_csv(self, file_path: Path) -> Optional[str]:
        """Convert a CSV into a Markdown table."""
        try:
            with open(file_path, "r", encoding="utf-8", newline="") as f:
                reader = csv.reader(f)
                rows = list(reader)
        except Exception as exc:
            print(f"[Warning] Failed to read CSV {file_path}: {exc}")
            return None
        if not rows:
            return ""
        header = rows[0]
        body = rows[1:]
        # Construct Markdown table
        table_lines: List[str] = []
        table_lines.append("|" + "|".join(h.strip() or " " for h in header) + "|")
        table_lines.append("|" + "|".join(["---"] * len(header)) + "|")
        for row in body:
            # Ensure row length matches header length
            padded = row + [""] * (len(header) - len(row))
            table_lines.append("|" + "|".join(cell.strip() or " " for cell in padded) + "|")
        return "\n".join(table_lines).strip()

    def _process_json(self, file_path: Path) -> Optional[str]:
        """Serialize JSON to a pretty‑printed string inside a code block."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception as exc:
            print(f"[Warning] Failed to read JSON {file_path}: {exc}")
            return None
        pretty = json.dumps(data, indent=2, ensure_ascii=False)
        # Wrap in a fenced code block so that markdown parsers treat it as
        # code.  This helps prevent headings or list markers inside the
        # JSON from being interpreted as markdown syntax.
        return f"```json\n{pretty}\n```"

    def _process_pptx(self, file_path: Path) -> Optional[str]:
        """Extract text from a PowerPoint file.

        If python‑pptx is not available, return ``None``.
        """
        if Presentation is None:
            print(f"[Warning] python-pptx is not installed; skipping PPTX {file_path}")
            return None
        try:
            pres = Presentation(str(file_path))  # type: ignore[arg-type]
        except Exception as exc:
            print(f"[Warning] Failed to open PPTX {file_path}: {exc}")
            return None
        lines: List[str] = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(getattr(shape, "text", "")).strip()
                    if text:
                        lines.append(text)
                elif hasattr(shape, "text_frame"):
                    try:
                        frame = shape.text_frame  # type: ignore[attr-defined]
                        for p in frame.paragraphs:
                            txt = p.text.strip()
                            if txt:
                                lines.append(txt)
                    except Exception:
                        continue
        return "\n".join(lines).strip()